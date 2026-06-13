from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 32, 65)  # Dark blue
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(100, 200, 255)
    
    # Presenter info
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    info_frame = info_box.text_frame
    info_frame.word_wrap = True
    
    p = info_frame.paragraphs[0]
    p.text = "Paridhi Jain"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(200, 200, 200)
    
    p = info_frame.add_paragraph()
    p.text = "HackArena 2 - Solo Entry"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(150, 150, 150)

def add_content_slide(prs, title, content_points):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 247, 250)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(15, 32, 65)
    
    # Line under title
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(4), Inches(0))
    line.line.color.rgb = RGBColor(0, 102, 204)
    line.line.width = Pt(3)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(8.6), Inches(4.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.level = 0
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_before = Pt(8)
        p.space_after = Pt(8)

def add_solution_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 247, 250)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Solution: What is DIP Arena?"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(15, 32, 65)
    
    # Line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(5), Inches(0))
    line.line.color.rgb = RGBColor(0, 102, 204)
    line.line.width = Pt(3)
    
    # Main text
    main_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(8.6), Inches(0.8))
    main_frame = main_box.text_frame
    main_frame.text = "Multi-Agent Decision Intelligence Platform"
    main_frame.paragraphs[0].font.size = Pt(26)
    main_frame.paragraphs[0].font.bold = True
    main_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    
    # Features
    features_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.8), Inches(8), Inches(3.5))
    features_frame = features_box.text_frame
    features_frame.word_wrap = True
    
    features = [
        "Multi-agent AI debate with structured reasoning",
        "Risk-aware recommendations",
        "Validation of claims with evidence scoring",
        "Transparent decision analysis",
        "Comparison between decision reports"
    ]
    
    for i, feature in enumerate(features):
        if i == 0:
            p = features_frame.paragraphs[0]
        else:
            p = features_frame.add_paragraph()
        p.text = feature
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_before = Pt(10)
        p.level = 0

def add_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 247, 250)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Architecture & Workflow"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(15, 32, 65)
    
    # Line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(4), Inches(0))
    line.line.color.rgb = RGBColor(0, 102, 204)
    line.line.width = Pt(3)
    
    # Workflow text
    workflow_text = """User Query
        ↓
Agent Orchestrator
        ↓
Economics  Ethics  Policy
Research  Risk  Strategy
        ↓
Validation Engine
        ↓
Decision Report
        ↓
Comparison + History"""
    
    workflow_box = slide.shapes.add_textbox(Inches(2), Inches(1.7), Inches(6), Inches(4.5))
    workflow_frame = workflow_box.text_frame
    workflow_frame.text = workflow_text
    workflow_frame.paragraphs[0].font.size = Pt(18)
    workflow_frame.paragraphs[0].font.color.rgb = RGBColor(40, 40, 40)
    workflow_frame.paragraphs[0].font.bold = True
    workflow_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_features_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 247, 250)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Key Features"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(15, 32, 65)
    
    # Line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(3), Inches(0))
    line.line.color.rgb = RGBColor(0, 102, 204)
    line.line.width = Pt(3)
    
    features = [
        "Multi-agent debate engine",
        "Validation matrix with evidence scoring",
        "Risk and tradeoff analysis",
        "Confidence scoring",
        "Counterfactual reasoning",
        "Report comparison",
        "Query history tracking",
        "Audit-friendly outputs"
    ]
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(4.5), Inches(4.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, feature in enumerate(features[:4]):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = "• " + feature
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_before = Pt(6)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.5), Inches(4.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for i, feature in enumerate(features[4:]):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = "• " + feature
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_before = Pt(6)

def add_innovation_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 247, 250)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Why It's Different"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(15, 32, 65)
    
    # Line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(4), Inches(0))
    line.line.color.rgb = RGBColor(0, 102, 204)
    line.line.width = Pt(3)
    
    innovations = [
        "Multi-perspective reasoning (not single-model AI)",
        "Explicit dissent instead of hidden reasoning",
        "Transparent AI governance approach",
        "Structured decision intelligence framework",
        "Counterfactual analysis capabilities",
        "Validation-aware recommendations"
    ]
    
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(8.6), Inches(4.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, innovation in enumerate(innovations):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        p.text = "✓ " + innovation
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0, 102, 204)
        p.font.bold = True
        p.space_before = Pt(10)

def add_usecases_slide(prs):
    add_content_slide(prs, "Real-World Use Cases", [
        "Public policy decisions with multi-stakeholder validation",
        "Healthcare resource allocation and triage",
        "Disaster response and emergency planning",
        "Regulatory governance and compliance",
        "Enterprise strategic analysis",
        "Risk assessment systems"
    ])

def add_techstack_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 247, 250)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Technology Stack"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(15, 32, 65)
    
    # Line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(3), Inches(0))
    line.line.color.rgb = RGBColor(0, 102, 204)
    line.line.width = Pt(3)
    
    tech = [
        "Backend: Python with ThreadPoolExecutor",
        "Multi-agent orchestration engine",
        "Structured reasoning framework",
        "Frontend: HTML5, CSS3, JavaScript",
        "Local-first architecture (no external API required)",
        "JSON-based report export"
    ]
    
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(8.6), Inches(4.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, item in enumerate(tech):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        p.text = "⚙ " + item
        p.font.size = Pt(19)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_before = Pt(10)

def add_future_slide(prs):
    add_content_slide(prs, "Future Roadmap", [
        "Live LLM integration for dynamic agent responses",
        "Real-time evidence retrieval from knowledge bases",
        "Explainable AI governance layer",
        "Collaborative decision workspace",
        "Enterprise deployment support",
        "Automated PDF report export"
    ])

def add_closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 32, 65)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Thank You!"
    title_frame.paragraphs[0].font.size = Pt(60)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    contact_frame = contact_box.text_frame
    contact_frame.word_wrap = True
    
    p = contact_frame.paragraphs[0]
    p.text = "GitHub: github.com/paridhijain5/dip-arena"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(100, 200, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p = contact_frame.add_paragraph()
    p.text = "Demo: http://127.0.0.1:8000"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(100, 200, 255)
    p.alignment = PP_ALIGN.CENTER

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Add all slides
add_title_slide(prs, "DIP Arena", "Multi-Agent Decision Intelligence Platform\nAI-powered structured reasoning for high-stakes decisions")
add_content_slide(prs, "The Problem", [
    "High-stakes decisions are often fragmented across teams",
    "Single-model AI responses lack transparency and auditability",
    "No structured disagreement or validation mechanisms",
    "Difficult to track reasoning or audit AI recommendations",
    "",
    "Critical decisions need structured reasoning, not just generated answers."
])
add_solution_slide(prs)
add_architecture_slide(prs)
add_features_slide(prs)
add_innovation_slide(prs)
add_usecases_slide(prs)
add_techstack_slide(prs)
add_future_slide(prs)
add_closing_slide(prs)

# Save
prs.save("c:/New folder (2)/DIP_Arena_Presentation.pptx")
print("[+] Presentation created: DIP_Arena_Presentation.pptx")
